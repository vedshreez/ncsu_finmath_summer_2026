#!/usr/bin/env python
"""Generate per-week student presentations (PPTX + Beamer PDF) for the
10-week program. One shared content source; PPTX via python-pptx, PDF via
LaTeX Beamer. NC State theme (Wolfpack Red, NCSU wordmark)."""
import os, subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE = "/Users/dendisuhubdy/Github/ncsu"
OUT  = os.path.join(HERE, "weekly_presentations")
LOGO = os.path.join(HERE, "assets", "ncsu_logo.png")
os.makedirs(OUT, exist_ok=True)

RED   = RGBColor(0xCC, 0x00, 0x00)
DKRED = RGBColor(0x99, 0x00, 0x00)
GRAY  = RGBColor(0x6D, 0x6D, 0x6D)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Arial"
SW, SH = Inches(13.333), Inches(7.5)

MENTOR = "Industry Mentor: Dendi Suhubdy (Bitwyre)  ·  Guest: Fenni Kang (Coincall)"

# Each week: (filename_stem, week_label, title, subtitle, [ (slide_head, [bullets]) ])
WEEKS = [
("week01_introduction", "Week 1", "Introduction",
 "Program overview, tooling, and data",
 [("Welcome to the Program",
   ["A 10-week summer research program on optimal market making for crypto options",
    "Goal: design, implement, and evaluate quoting strategies on real Coincall data",
    "You will produce a research-grade report and a working code repository",
    "This week: orientation, environment setup, and data access"]),
  ("Maker vs. Taker",
   ["A market maker continuously quotes both sides and earns the spread",
    "A taker crosses the spread to execute immediately",
    "We build the maker's side; in Week 9 a Coincall practitioner covers taker strategies",
    "The spread is compensation for committing to quotes under uncertainty"]),
  ("The Spine: Avellaneda-Stoikov",
   ["The whole course is organized around the Avellaneda-Stoikov framework",
    "Introduced Week 3; turned into a calibrated engine Week 5",
    "Extended to multi-Greek options inventory Week 8",
    "Advanced extensions (multi-asset, adverse selection, robust) in Week 10"]),
  ("Tooling",
   ["Problem sets: Python + PyTorch (autograd does the differentiation for you)",
    "You will NOT hand-roll algorithmic differentiation",
    "Fast-computation week (Week 7): C++ exposed to Python via pybind11",
    "Industry pattern: Python for research, C++ for the hot path"]),
  ("Data, Milestones, Assessment",
   ["Coincall BTC/ETH order books, trades, options chains, funding rates",
    "Milestone 1 (Wk5): Avellaneda-Stoikov engine",
    "Milestone 2 (Wk7): fast pricing & Greeks library",
    "Milestone 3 (Wk9): full options market-making backtest",
    "Grade: problem sets 30%, milestones 30%, final 30%, participation 10%"]),
  ("This Week's To-Do",
   ["Set up Python + PyTorch and a C++17 toolchain with CMake",
    "Build and import the pybind11 'hello world' starter module",
    "Get your Coincall data credentials and sign the data-use agreement",
    "No graded problem set this week"])]),

("week02_microstructure", "Week 2", "Foundations of Market Microstructure",
 "From limit order books to the role of the market maker",
 [("Learning Objectives",
   ["Reconstruct a limit order book from raw exchange events",
    "Decompose the bid-ask spread into its three classical components",
    "Explain the market maker as a continuous liquidity provider"]),
  ("Lecture 1: Order Book Mechanics",
   ["Order types: market, limit, IOC, post-only, hidden",
    "Price-time priority and the matching algorithm",
    "L1 / L2 / L3 views; tick size and lot size",
    "Coincall architecture: contracts, fees, market-maker tiers"]),
  ("Lecture 2: Spread Decomposition",
   ["Adverse selection, inventory, and order-processing costs",
    "Roll (1984): effective spread from trade prices",
    "Glosten-Milgrom (1985): asymmetric information",
    "Kyle (1985): the strategic informed trader",
    "Realized spread, effective spread, price impact"]),
  ("Recitation",
   ["Build an L2 order-book reconstructor in Python",
    "Replay one hour of Coincall BTC-PERP events",
    "Emit a CSV of mid-price, spread, and depth",
    "Watch out: delete price levels that empty after a cancel"]),
  ("Reading & Problem Set 2 (Python)",
   ["Cartea-Jaikumar-Penalva Ch. 1-3; Roll (1984); Glosten-Milgrom (1985)",
    "PS2: order book, Roll estimator, Glosten-Milgrom, toxicity"])]),

("week03_avellaneda_stoikov", "Week 3", "The Avellaneda-Stoikov Model",
 "A worked example of optimal market making via stochastic control",
 [("Learning Objectives",
   ["Set up the Avellaneda-Stoikov control problem with CARA utility",
    "Understand the HJB equation and its reduction via the ansatz",
    "Interpret the reservation price and optimal half-spread"]),
  ("Lecture 1: Setup and Dynamics",
   ["Mid-price as arithmetic Brownian motion (chosen for tractability)",
    "Poisson order flow with intensity lambda(d) = A exp(-kappa d)",
    "Wealth and inventory dynamics under continuous quoting",
    "Exponential (CARA) utility decouples value from wealth"]),
  ("Lecture 2: HJB and Closed-Form Quotes",
   ["Dynamic programming principle -> HJB equation",
    "Separation ansatz V = -exp(-gamma(x+qS)) theta(t,q)",
    "Reservation price r = S - q gamma sigma^2 (T-t)",
    "Half-spread = (1/gamma) ln(1+gamma/kappa) + (gamma sigma^2 / 2)(T-t)"]),
  ("Recitation",
   ["Implement the closed-form quotes in PyTorch",
    "Visualize reservation price and spread vs. inventory",
    "Run a risk-aversion (gamma) sensitivity sweep"]),
  ("Reading & Problem Set 3 (Python/PyTorch)",
   ["Avellaneda-Stoikov (2008); Cartea-Jaikumar-Penalva Ch. 10",
    "PS3: ODE solver, closed-form quotes, simulation, gamma sweep"])]),

("week04_control", "Week 4", "Stochastic Optimal Control in Continuous Time",
 "The general theory behind the worked example",
 [("Learning Objectives",
   ["Apply the dynamic programming principle in continuous time",
    "Understand the HJB equation and the verification theorem",
    "See why CARA and CRRA utilities give tractable problems"]),
  ("Lecture 1: Dynamic Programming & HJB",
   ["Controlled diffusions and admissible controls",
    "The dynamic programming principle",
    "The infinitesimal generator",
    "Deriving the HJB equation; boundary/terminal conditions"]),
  ("Lecture 2: Verification & Utilities",
   ["The verification theorem: a smooth HJB solution IS the value function",
    "Viscosity solutions: when value functions are non-smooth",
    "CARA vs. CRRA vs. HARA",
    "Why CARA decouples the value function from wealth"]),
  ("Recitation",
   ["Solve the Merton problem under CARA, then CRRA",
    "Compare the two solutions side by side",
    "Recognize the recurring HJB -> ansatz -> ODE -> solution pattern"]),
  ("Reading & Problem Set 4 (Python/PyTorch)",
   ["Pham (2009) Ch. 3-4; Cartea-Jaikumar-Penalva Ch. 6",
    "PS4: Merton CRRA/CARA, Almgren-Chriss, HJB-residual check"])]),

("week05_engine_milestone1", "Week 5", "The Avellaneda-Stoikov Engine",
 "Implementation, calibration, and backtest -- Milestone 1",
 [("Learning Objectives",
   ["Implement the quotes as a clean, tested PyTorch engine",
    "Calibrate the order-flow intensities by maximum likelihood",
    "Build an event-driven backtest and benchmark on Coincall data"]),
  ("Lecture 1: From Closed Form to Code",
   ["Engine architecture; state management from the L2 book",
    "Mapping continuous-time quantities to discrete events",
    "Calibrating lambda(d) = A exp(-kappa d) from fill data",
    "Pitfalls: units, cash-process sign, inventory indexing"]),
  ("Lecture 2: Calibration & Validation",
   ["Maximum-likelihood fit of A and kappa with goodness-of-fit",
    "The backtest harness: replay, fill simulator, P&L attribution",
    "Validate the P&L distribution against theory",
    "Reproducibility under a fixed seed"]),
  ("Milestone 1 Deliverable",
   ["Python/PyTorch Avellaneda-Stoikov package on one linear asset",
    "Calibration + simulation + one-month BTC-PERP backtest",
    "Acceptance: quotes match closed form; P&L within +/-50% of reference",
    "pytest coverage; report <= 10 pages"])]),

("week06_vol_surface", "Week 6", "Options Pricing, the Volatility Surface, and Calibration",
 "Building a volatility surface from messy crypto data",
 [("Learning Objectives",
   ["Recall the Greeks, including vanna and volga",
    "State the no-arbitrage constraints on the vol surface",
    "Calibrate SVI and SABR with attention to stability"]),
  ("Lecture 1: Greeks & the Surface",
   ["Black-Scholes and the Greek vector",
    "Vanna, volga, and volatility risk",
    "The implied-vol smile and term structure",
    "Calendar and butterfly arbitrage"]),
  ("Lecture 2: SVI & SABR Calibration",
   ["Raw SVI and the arbitrage-free formulation",
    "SVI calibration by nonlinear least squares",
    "The SABR formula and the Hagan expansion",
    "Time-stability of calibrated parameters"]),
  ("Recitation & Problem Set 6 (Python/PyTorch)",
   ["Live calibration in PyTorch (autograd Jacobians)",
    "PS6: BS Greeks via autograd, SVI/SABR fits, stability study",
    "Reading: Gatheral (2006); Gatheral-Jacquier (2014); Hagan et al."])]),

("week07_fast_computation_milestone2", "Week 7",
 "Fast Computation of Prices, Greeks, Implied Vol, and Surfaces",
 "The latency budget, in C++ via pybind11 -- Milestone 2",
 [("Learning Objectives",
   ["Implement characteristic-function pricers in C++ (COS, Carr-Madan)",
    "Implement robust implied-vol inversion and surface revaluation",
    "Expose the C++ kernels to Python with pybind11"]),
  ("Lecture 1: CF Pricing & Implied Vol in C++",
   ["Characteristic functions under Black-Scholes and Heston",
    "The Fang-Oosterlee COS method (the workhorse)",
    "Robust implied-vol inversion: good guess + safeguarded Newton",
    "Vectorizing a whole surface"]),
  ("Lecture 2: Binding C++ to Python with pybind11",
   ["How a pybind11 module is structured and built with CMake",
    "Passing NumPy/PyTorch tensors with no copy (buffer protocol)",
    "Releasing the GIL around the C++ hot loop",
    "PyTorch autograd is the REFERENCE for Greeks; C++ is the speed"]),
  ("Milestone 2 Deliverable",
   ["C++ pricing/implied-vol/Greeks library, imported via pybind11",
    "Full BTC surface (300-500 contracts) revalued under 5 ms",
    "Greeks within 1e-4 of PyTorch autograd; pricing within 1e-4 of reference",
    "Latency reported as a distribution, plus a BINDING.md explainer"])]),

("week08_options_mm_hedging", "Week 8", "Options Market Making Theory and Dynamic Hedging",
 "Avellaneda-Stoikov with a vector-valued inventory",
 [("Learning Objectives",
   ["Set up the El Aoud-Abergel multi-Greek framework",
    "Understand the inventory penalty as a quadratic form",
    "Quantify the P&L attribution of a hedged options book"]),
  ("Lecture 1: Multi-Greek Inventory & HJB",
   ["Inventory becomes a vector of Greeks (delta, gamma, vega, ...)",
    "Quadratic inventory penalty and its economic meaning",
    "The multi-Greek HJB and separation ansatz",
    "Reduces to Avellaneda-Stoikov in the 1-D case"]),
  ("Lecture 2: Dynamic Hedging",
   ["Continuous and discrete delta hedging",
    "The gamma-theta-variance identity (the key to options trading)",
    "External vs. integrated hedging architectures",
    "Perpetual futures as the hedge instrument; funding cost"]),
  ("Recitation & Problem Set 8 (Python/PyTorch)",
   ["Multi-Greek HJB walk-through; PyTorch hedged-book simulation",
    "PS8: inventory dynamics, penalty calibration, threshold hedger",
    "Reading: El Aoud-Abergel (2015); Taleb (1997)"])]),

("week09_backtesting_taker_milestone3", "Week 9",
 "Backtesting, Risk Analysis, and Taker Strategies",
 "Assembling the full system -- Milestone 3 (with a Coincall guest session)",
 [("Learning Objectives",
   ["Build an event-driven backtest for options market making",
    "Produce a full P&L attribution and risk metrics",
    "Understand taker strategies and order-flow toxicity"]),
  ("Lecture 1: Backtest Architecture & Risk",
   ["Event-driven simulation with deterministic replay",
    "The fill simulator; state tracking of inventory and Greeks",
    "Pitfalls: look-ahead bias, unrealistic fills, expiry handling",
    "Sharpe/Calmar, drawdown, block-bootstrap confidence intervals"]),
  ("Lecture 2 (Guest: Fenni Kang, Coincall): Taker Strategies",
   ["When and why takers cross the spread",
    "Signals behind aggressive flow: momentum, liquidations, news",
    "Execution styles: sweep, TWAP/VWAP, iceberg-taking",
    "What taker behaviour means for the maker's order-flow toxicity"]),
  ("Milestone 3 Deliverable",
   ["Full multi-Greek backtest on 3 months of BTC & ETH options",
    "P&L attribution sums to total within 1 bp; Sharpe with CI",
    "Greek exposures bounded; drawdown forensic + taker analysis",
    "Reproducible from config; report <= 20 pages"])]),

("week10_advanced_structured_products", "Week 10",
 "Advanced Avellaneda-Stoikov, Structured Products, and Final Presentations",
 "Extensions, the corporate view, and communicating results (Coincall guest session)",
 [("Learning Objectives",
   ["Implement one advanced Avellaneda-Stoikov extension",
    "Understand how corporates use options via structured products",
    "Write and present a research-grade report"]),
  ("Lecture 1: Advanced Extensions",
   ["Multi-asset market making across BTC and ETH options",
    "Adverse-selection-aware and signal-driven quoting",
    "Discrete-inventory formulations",
    "Robust optimization with a worst-case inventory penalty"]),
  ("Lecture 2 (Guest: Fenni Kang, Coincall): Structured Products",
   ["Collars to fix a price band on a holding",
    "Accumulators / decumulators for staged accumulation or liquidation",
    "Coupon / yield-enhancement products",
    "How structured-product flow lands on the market-maker's book"]),
  ("Final Deliverables",
   ["One-page extension proposal (day one)",
    "PyTorch implementation on the milestone infrastructure",
    "Quantitative comparison vs. the Week-9 baseline",
    "Final report (<= 30 pages) and 20-minute presentation"])]),
]


# ----------------------------------------------------------------- PPTX
def _set(run, size, color, bold=False, italic=False):
    run.font.name = FONT; run.font.size = Pt(size)
    run.font.color.rgb = color; run.font.bold = bold; run.font.italic = italic


def build_pptx(week):
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    blank = prs.slide_layouts[6]
    # title slide
    s = prs.slides.add_slide(blank)
    band = s.shapes.add_shape(1, 0, 0, SW, Inches(2.5))
    band.fill.solid(); band.fill.fore_color.rgb = RED; band.line.fill.background()
    s.shapes.add_picture(LOGO, Inches(0.5), Inches(0.5), width=Inches(3.0))
    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(12.1), Inches(2.4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = week["label"] + " — " + week["title"]
    _set(r, 32, DKRED, bold=True)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = week["subtitle"]; _set(r2, 19, GRAY, italic=True)
    mb = s.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.8))
    mp = mb.text_frame.paragraphs[0]; mr = mp.add_run()
    mr.text = "Optimal Market Making for Cryptocurrency Options  ·  NC State University"
    _set(mr, 13, BLACK)
    rule = s.shapes.add_shape(1, Inches(0.6), Inches(5.9), Inches(4.0), Pt(3))
    rule.fill.solid(); rule.fill.fore_color.rgb = RED; rule.line.fill.background()
    # content slides
    for head, bullets in week["slides"]:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(LOGO, Inches(10.7), Inches(0.3), width=Inches(2.2))
        hb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.8), Inches(0.9))
        hp = hb.text_frame.paragraphs[0]; hr = hp.add_run(); hr.text = head; _set(hr, 26, DKRED, bold=True)
        rule = s.shapes.add_shape(1, Inches(0.5), Inches(1.35), Inches(12.3), Pt(2.5))
        rule.fill.solid(); rule.fill.fore_color.rgb = RED; rule.line.fill.background()
        bb = s.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.4))
        btf = bb.text_frame; btf.word_wrap = True
        for i, b in enumerate(bullets):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            run = p.add_run(); run.text = "•  " + b; _set(run, 19, BLACK)
            p.space_after = Pt(10)
    path = os.path.join(OUT, week["stem"] + ".pptx")
    prs.save(path); return path


# ----------------------------------------------------------------- BEAMER
def esc(t):
    for a, b in [("\\", r"\textbackslash{}"), ("&", r"\&"), ("%", r"\%"),
                 ("#", r"\#"), ("_", r"\_"), ("{", r"\{"), ("}", r"\}"),
                 ("$", r"\$"), ("~", r"\textasciitilde{}"), ("^", r"\textasciicircum{}"),
                 ("<", r"\textless{}"), (">", r"\textgreater{}")]:
        t = t.replace(a, b)
    return t

BEAMER_PRE = r"""\documentclass[aspectratio=169,11pt]{beamer}
\usetheme{default}
\setbeamertemplate{navigation symbols}{}
\usepackage{graphicx}
\usepackage[T1]{fontenc}
\definecolor{ncsured}{HTML}{CC0000}
\definecolor{ncsudk}{HTML}{990000}
\setbeamercolor{structure}{fg=ncsured}
\setbeamercolor{frametitle}{fg=white,bg=ncsured}
\setbeamercolor{title}{fg=ncsudk}
\setbeamercolor{itemize item}{fg=ncsured}
\setbeamercolor{itemize subitem}{fg=ncsured}
\setbeamerfont{frametitle}{series=\bfseries}
\setbeamertemplate{frametitle}{%
  \nointerlineskip\vskip2pt%
  \begin{beamercolorbox}[wd=\paperwidth,ht=2.4ex,dp=1ex,leftskip=0.3cm]{frametitle}%
  \bfseries\insertframetitle%
  \end{beamercolorbox}}
\setbeamertemplate{footline}{%
  \hbox{\begin{beamercolorbox}[wd=\paperwidth,ht=2.2ex,dp=1ex,leftskip=.3cm,rightskip=.3cm]{}%
  \tiny Optimal Market Making for Cryptocurrency Options \hfill __LABEL__ \hfill \insertframenumber%
  \end{beamercolorbox}}}
"""

def build_beamer_tex(week):
    body = [BEAMER_PRE.replace("__LABEL__", esc(week["label"]))]
    body.append(r"\title{%s: %s}" % (esc(week["label"]), esc(week["title"])))
    body.append(r"\subtitle{%s}" % esc(week["subtitle"]))
    body.append(r"\date{North Carolina State University \,\textbar\, Summer 2026}")
    body.append(r"\titlegraphic{\includegraphics[width=4.2cm]{%s}}" % LOGO)
    body.append(r"\begin{document}")
    # title frame
    body.append(r"\begin{frame}[plain]\titlepage")
    body.append(r"\begin{center}\tiny " + esc(MENTOR) + r"\end{center}\end{frame}")
    for head, bullets in week["slides"]:
        body.append(r"\begin{frame}{%s}" % esc(head))
        body.append(r"\begin{itemize}")
        for b in bullets:
            body.append(r"  \item %s" % esc(b))
        body.append(r"\end{itemize}\end{frame}")
    body.append(r"\end{document}")
    path = os.path.join(OUT, week["stem"] + ".tex")
    open(path, "w", encoding="utf-8").write("\n".join(body))
    return path


def to_dict(t):
    stem, label, title, subtitle, slides = t
    return {"stem": stem, "label": label, "title": title,
            "subtitle": subtitle, "slides": slides}


if __name__ == "__main__":
    for t in WEEKS:
        w = to_dict(t)
        build_pptx(w)
        build_beamer_tex(w)
        print("built", w["stem"])
    print("Done: %d weeks." % len(WEEKS))
