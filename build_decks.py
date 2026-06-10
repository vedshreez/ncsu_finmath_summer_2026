#!/usr/bin/env python
"""Build NC State-themed PPTX decks for each course document."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# NC State theme (from OptimalMarketMaking_WeeklyUpdate_06_08.pptx theme1.xml)
RED   = RGBColor(0xCC, 0x00, 0x00)   # Wolfpack Red (primary)
DKRED = RGBColor(0x99, 0x00, 0x00)
GRAY  = RGBColor(0x6D, 0x6D, 0x6D)
LGRAY = RGBColor(0xC3, 0xC3, 0xC3)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL  = RGBColor(0x00, 0x84, 0x73)   # NCSU secondary
FONT  = "Arial"
LOGO  = "/Users/dendisuhubdy/Github/ncsu/assets/ncsu_logo.png"

SW, SH = Inches(13.333), Inches(7.5)   # 16:9


def _set(run, size, color, bold=False, italic=False):
    run.font.name = FONT; run.font.size = Pt(size)
    run.font.color.rgb = color; run.font.bold = bold; run.font.italic = italic


def add_logo(slide, top=Inches(0.28), width=Inches(2.3)):
    slide.shapes.add_picture(LOGO, Inches(0.45), top, width=width)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_slide(prs, title, subtitle, meta_lines):
    s = blank(prs)
    # red band across top
    band = s.shapes.add_shape(1, 0, 0, SW, Inches(2.5))
    band.fill.solid(); band.fill.fore_color.rgb = RED
    band.line.fill.background()
    # logo (white-on-red) sits on the red band
    s.shapes.add_picture(LOGO, Inches(0.5), Inches(0.45), width=Inches(3.0))
    # title
    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.9), Inches(12.1), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title; _set(r, 34, DKRED, bold=True)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle; _set(r2, 20, GRAY, italic=True)
    # meta
    mb = s.shapes.add_textbox(Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.7))
    mtf = mb.text_frame; mtf.word_wrap = True
    for i, line in enumerate(meta_lines):
        p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
        r = p.add_run(); r.text = line; _set(r, 15, BLACK, bold=(i == 0))
    # red rule
    rule = s.shapes.add_shape(1, Inches(0.6), Inches(5.25), Inches(4.0), Pt(3))
    rule.fill.solid(); rule.fill.fore_color.rgb = RED; rule.line.fill.background()
    return s


def content_slide(prs, heading, bullets, sub=None):
    s = blank(prs)
    add_logo(s)
    # heading
    hb = s.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.9))
    htf = hb.text_frame; htf.word_wrap = True
    p = htf.paragraphs[0]; r = p.add_run(); r.text = heading; _set(r, 26, DKRED, bold=True)
    if sub:
        sp = htf.add_paragraph(); sr = sp.add_run(); sr.text = sub; _set(sr, 14, GRAY, italic=True)
    # red rule under heading
    rule = s.shapes.add_shape(1, Inches(0.5), Inches(1.78), Inches(12.3), Pt(2.5))
    rule.fill.solid(); rule.fill.fore_color.rgb = RED; rule.line.fill.background()
    # bullets
    bb = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.1), Inches(5.1))
    btf = bb.text_frame; btf.word_wrap = True
    for i, (lvl, text, *style) in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.level = lvl
        r = p.add_run()
        bullet = ("•  " if lvl == 0 else "–  ")
        r.text = bullet + text
        color = RED if (style and style[0] == "accent") else BLACK
        bold = bool(style and style[0] == "accent")
        _set(r, 18 if lvl == 0 else 16, color, bold=bold)
        p.space_after = Pt(6)
    return s


def section_slide(prs, label, title):
    s = blank(prs)
    band = s.shapes.add_shape(1, 0, Inches(2.6), SW, Inches(2.3))
    band.fill.solid(); band.fill.fore_color.rgb = RED; band.line.fill.background()
    add_logo(s)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.95), Inches(12), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = label; _set(r, 16, WHITE, bold=True)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = title; _set(r2, 30, WHITE, bold=True)
    return s


def save(prs, path):
    prs.save(path); print("wrote", path)


def new_prs():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    return prs


META = ["North Carolina State University — Department of Mathematics",
        "Industry Mentor: Dendi Suhubdy (Bitwyre, San Francisco)",
        "Guest Practitioner: Fenni Kang (Chief Strategy Officer, Coincall)",
        "Summer 2026"]

# ---------------------------------------------------------------- SYLLABUS
def build_syllabus():
    prs = new_prs()
    title_slide(prs, "Optimal Market Making for Cryptocurrency Options",
                "The Avellaneda–Stoikov Framework and its Extension to Options — A 10-Week Program",
                META)
    content_slide(prs, "Course Overview", [
        (0, "Re-centered on the Avellaneda–Stoikov framework for options market making", "accent"),
        (0, "Week 1 was the introduction (delivered last week); Weeks 2–10 are technical"),
        (0, "Extend single-asset Avellaneda–Stoikov to multi-Greek options inventory"),
        (0, "Live and historical Coincall data throughout; research-grade final report"),
        (0, "Advanced week covers Avellaneda–Stoikov extensions (no GLFT)"),
    ])
    content_slide(prs, "Implementation Languages", [
        (0, "Problem sets in Python, using PyTorch", "accent"),
        (1, "Differentiation is implicit via autograd — no hand-rolled AD"),
        (1, "Greeks, surface calibration, and intensity fits all via .backward()"),
        (0, "Fast-computation methods in C++", "accent"),
        (1, "Latency-critical pricers, implied-vol, and surface revaluation"),
        (1, "Exposed to Python through a pybind11 binding (taught in Week 7)"),
        (0, "Industry pattern: Python for research, C++ for the hot path, pybind11 between"),
    ])
    content_slide(prs, "Ten-Week Schedule", [
        (0, "Wk 1 — Introduction (delivered)"),
        (0, "Wk 2 — Market Microstructure"),
        (0, "Wk 3 — The Avellaneda–Stoikov Model"),
        (0, "Wk 4 — Stochastic Optimal Control"),
        (0, "Wk 5 — Avellaneda–Stoikov Engine  ·  Milestone 1", "accent"),
        (0, "Wk 6 — Volatility Surface & Calibration"),
        (0, "Wk 7 — Fast Computation: prices, Greeks, implied vol, surfaces  ·  Milestone 2", "accent"),
        (0, "Wk 8 — Options MM Theory & Dynamic Hedging"),
        (0, "Wk 9 — Backtesting, Risk & Taker Strategies  ·  Milestone 3", "accent"),
        (0, "Wk 10 — Advanced Avellaneda–Stoikov, Structured Products & Finals"),
    ])
    content_slide(prs, "Milestones", [
        (0, "Milestone 1 (Wk 5) — Linear-asset Avellaneda–Stoikov engine (Python/PyTorch)"),
        (1, "Calibrated, backtested on one month of Coincall BTC-PERP"),
        (0, "Milestone 2 (Wk 7) — Fast pricing & Greeks library (C++ via pybind11)"),
        (1, "Full BTC surface revalued < 5 ms; Greeks validated vs. PyTorch autograd"),
        (0, "Milestone 3 (Wk 9) — Full options market-making backtest"),
        (1, "Three months BTC & ETH; P&L attribution, risk metrics, Greek exposures"),
    ])
    content_slide(prs, "Guest Practitioner — Fenni Kang (Coincall)",
        [
        (0, "Week 9 — Taker Strategies", "accent"),
        (1, "When/why takers cross the spread; signals behind aggressive flow"),
        (1, "What taker behaviour implies for order-flow toxicity"),
        (0, "Week 10 — Corporate Usage of Options & Structured Products", "accent"),
        (1, "Collars to fix a price band"),
        (1, "Accumulators / decumulators for staged accumulation or liquidation"),
        (1, "Coupon / yield-enhancement products"),
        (1, "How structured-product flow lands on the market-maker's book"),
    ])
    content_slide(prs, "Assessment", [
        (0, "Weekly problem sets — 30%"),
        (0, "Three milestone deliverables (Wk 5, 7, 9) — 30%"),
        (0, "Final report and presentation — 30%"),
        (0, "Participation in lecture, recitation, code review — 10%"),
    ])
    save(prs, "/Users/dendisuhubdy/Github/ncsu/syllabus.pptx")

# ---------------------------------------------------------------- HANDBOOK
def build_handbook():
    prs = new_prs()
    title_slide(prs, "Tutor Handbook",
                "Optimal Market Making for Cryptocurrency Options — Avellaneda–Stoikov for Options",
                META)
    content_slide(prs, "Your Role", [
        (0, "Lead the weekly 90-minute recitation"),
        (0, "Hold three office hours per week"),
        (0, "Grade problem sets within seven days"),
        (0, "First-pass milestone evaluation (mentor gives the final grade)"),
        (0, "Flag at-risk students to the mentor within 24 hours"),
        (0, "Two lectures are guest sessions (Fenni Kang, Coincall) — Wk 9 & Wk 10"),
    ])
    content_slide(prs, "Grading Standards — Two Languages", [
        (0, "Code graded on correctness, reproducibility, quality — in that order", "accent"),
        (0, "Problem sets are Python/PyTorch; autograd is the oracle for Greeks"),
        (0, "Students are NOT asked to hand-derive algorithmic differentiation"),
        (0, "Week 7 also grades the C++/pybind11 boundary", "accent"),
        (1, "Builds with CMake, imports cleanly, agrees with PyTorch reference Greeks"),
        (0, "Hand-rolled adjoint AD in C++ is a red flag — steer back to autograd"),
    ])
    content_slide(prs, "Week-by-Week Spine", [
        (0, "Wk 2 Microstructure — read an order book; spread decomposition"),
        (0, "Wk 3 Avellaneda–Stoikov — the full HJB→ansatz→ODE→quote pipeline"),
        (0, "Wk 4 Control — the abstract template; verification theorem"),
        (0, "Wk 5 Engine (M1) — closed form becomes calibrated, backtested code"),
        (0, "Wk 6 Vol surface — SVI/SABR; calibration is half craft, half theory"),
        (0, "Wk 7 Fast computation (M2) — latency; the pybind11 binder lecture"),
        (0, "Wk 8 Options MM — vector inventory; gamma–theta–variance identity"),
        (0, "Wk 9 Backtesting (M3) — engineering discipline; taker flow guest"),
        (0, "Wk 10 Advanced AV + structured-products guest + finals"),
    ])
    content_slide(prs, "Milestone Rubrics (out of 100)", [
        (0, "M1 — correctness 30, calibration 20, reproducibility 15, backtest 15,"),
        (1, "vectorization 10, code quality 5, report 5"),
        (0, "M2 — pricing 20, Greeks 15, latency 25 (hard gate, caps at 70),", "accent"),
        (1, "pybind11 integration 10, implied-vol 10, tests 10, note 10"),
        (0, "M3 — reproducibility 10, attribution 20, risk 15, Greeks 10,"),
        (1, "engine 15, hedging 10, forensic 10, report 10"),
    ])
    content_slide(prs, "Escalate Immediately (⚑ phone, not Slack)", [
        (0, "Suspected academic-integrity violation"),
        (0, "Student in acute distress"),
        (0, "Data-security incident (e.g., Coincall data exposed publicly)"),
        (0, "Hostile interaction"),
    ])
    save(prs, "/Users/dendisuhubdy/Github/ncsu/handbook.pptx")

# ---------------------------------------------------------------- PROBLEM SETS
def build_problem_sets():
    prs = new_prs()
    title_slide(prs, "Problem Sets",
                "Optimal Market Making for Cryptocurrency Options — Weeks 2–10",
                META)
    content_slide(prs, "How the Problem Sets Work", [
        (0, "Python with PyTorch — autograd handles all derivatives", "accent"),
        (0, "Week 7 is C++, called from Python via pybind11 (scaffold provided)", "accent"),
        (0, "Each set ships with a starter repo, data loaders, and pytest harness"),
        (0, "Emit results as CSV and figures as PNG; set seeds for reproducibility"),
        (0, "Worked solutions distributed separately after each deadline"),
    ])
    content_slide(prs, "Weeks 2–5", [
        (0, "PS2 Microstructure — order book, Roll, Glosten–Milgrom, toxicity"),
        (0, "PS3 Avellaneda–Stoikov — ODE solver, closed-form quotes, γ sweep"),
        (0, "PS4 Control — Merton CRRA/CARA, Almgren–Chriss, HJB residual"),
        (0, "PS5 Engine (Milestone 1) — calibration, one-month backtest, tests", "accent"),
    ])
    content_slide(prs, "Weeks 6–10", [
        (0, "PS6 Vol surface — BS Greeks via autograd, SVI/SABR calibration"),
        (0, "PS7 Fast computation (Milestone 2) — C++ COS/implied-vol via pybind11", "accent"),
        (1, "Greeks validated against PyTorch autograd; surface < 5 ms"),
        (0, "PS8 Options MM — multi-Greek inventory, hedging integration"),
        (0, "PS9 Backtesting (Milestone 3) — attribution, drawdown + taker analysis", "accent"),
        (0, "PS10 Advanced extension + final report & presentation"),
    ])
    save(prs, "/Users/dendisuhubdy/Github/ncsu/problem_sets.pptx")

# ---------------------------------------------------------------- SOLUTIONS
def build_solutions():
    prs = new_prs()
    title_slide(prs, "Problem Set Solutions",
                "Confidential — graders, and students after each deadline",
                META)
    content_slide(prs, "Using These Solutions", [
        (0, "One correct approach per problem: sketch, expected result, grading notes"),
        (0, "Students may reach the same result by other means"),
        (0, "Grade on correctness, reproducibility, and code quality"),
        (0, "Point allocations align with the Handbook milestone rubrics"),
    ])
    content_slide(prs, "Recurring Grading Themes", [
        (0, "Order book: delete empty levels after a cancel (−2 if missed)"),
        (0, "Calibration: maximum likelihood, not OLS; report χ² goodness-of-fit"),
        (0, "Greeks: PyTorch autograd is the reference — verify against finite diff"),
        (0, "Week 7: latency is a hard gate; release the GIL around the C++ loop", "accent"),
        (0, "Attribution must sum to total P&L — a mismatch is always a bug", "accent"),
        (0, "Sharpe always with a block-bootstrap confidence interval"),
    ])
    content_slide(prs, "Milestone Solution Pointers", [
        (0, "M1 — quotes match closed form < 1e-6; backtest within ±50% of reference"),
        (0, "M2 — C++ surface < 5 ms; C++ Greeks within 1e-4 of autograd"),
        (0, "M3 — deterministic replay; attribution sums to total within 1 bp"),
    ])
    save(prs, "/Users/dendisuhubdy/Github/ncsu/solutions.pptx")


if __name__ == "__main__":
    build_syllabus(); build_handbook(); build_problem_sets(); build_solutions()
    print("All decks built.")
